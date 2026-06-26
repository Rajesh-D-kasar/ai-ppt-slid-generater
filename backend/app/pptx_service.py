from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .models import DeckPlan, SlidePlan


THEMES = {
    "modern": {
        "background": RGBColor(247, 250, 252),
        "surface": RGBColor(255, 255, 255),
        "primary": RGBColor(20, 83, 136),
        "accent": RGBColor(17, 145, 110),
        "muted": RGBColor(91, 111, 132),
        "text": RGBColor(24, 32, 43),
    },
    "dark": {
        "background": RGBColor(17, 24, 39),
        "surface": RGBColor(31, 41, 55),
        "primary": RGBColor(96, 165, 250),
        "accent": RGBColor(45, 212, 191),
        "muted": RGBColor(203, 213, 225),
        "text": RGBColor(248, 250, 252),
    },
    "warm": {
        "background": RGBColor(255, 251, 235),
        "surface": RGBColor(255, 255, 255),
        "primary": RGBColor(154, 52, 18),
        "accent": RGBColor(5, 150, 105),
        "muted": RGBColor(120, 113, 108),
        "text": RGBColor(41, 37, 36),
    },
    "minimal": {
        "background": RGBColor(255, 255, 255),
        "surface": RGBColor(248, 250, 252),
        "primary": RGBColor(15, 23, 42),
        "accent": RGBColor(37, 99, 235),
        "muted": RGBColor(100, 116, 139),
        "text": RGBColor(30, 41, 59),
    },
    "startup": {
        "background": RGBColor(250, 245, 255),
        "surface": RGBColor(255, 255, 255),
        "primary": RGBColor(88, 28, 135),
        "accent": RGBColor(234, 88, 12),
        "muted": RGBColor(87, 83, 78),
        "text": RGBColor(28, 25, 23),
    },
}

DECK_TYPE_LABELS = {
    "business": "Business Brief",
    "startup_pitch": "Startup Pitch",
    "education": "Education Deck",
    "sales": "Sales Deck",
    "research": "Research Report",
    "general": "General Deck",
}

SLIDE_W = 13.333
SLIDE_H = 7.5


def _set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text_box(slide, left, top, width, height, text=""):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if text:
        frame.text = text
    return shape


def _shape(slide, kind, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _add_run(paragraph, text, size, color, bold=False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def _style_existing_text(shape, size, color, bold=False, align=None):
    paragraph = shape.text_frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    if not paragraph.runs:
        paragraph.add_run()
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_footer(slide, theme: dict, slide_number: int, total: int) -> None:
    box = _text_box(slide, 0.65, 6.94, 12.05, 0.25)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = f"AI PPT Slide Generator  |  {slide_number}/{total}"
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph.runs[0].font.size = Pt(8)
    paragraph.runs[0].font.color.rgb = theme["muted"]


def _add_header(slide, title: str, theme: dict) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 1.2, 1.25, 0.06, theme["accent"])
    title_box = _text_box(slide, 0.65, 0.48, 12, 0.62)
    _add_run(title_box.text_frame.paragraphs[0], title, 29, theme["primary"], True)


def _add_notes_and_footer(slide, slide_plan: SlidePlan, theme: dict, number: int, total: int) -> None:
    if slide_plan.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_plan.speaker_notes
    _add_footer(slide, theme, number, total)


def _add_title_slide(presentation, deck: DeckPlan, theme: dict, deck_type: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, theme["background"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, SLIDE_H, theme["accent"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 4.52, 2.95, 0.08, theme["accent"])

    label = _text_box(slide, 0.78, 0.86, 4.8, 0.35, DECK_TYPE_LABELS.get(deck_type, "Presentation Deck"))
    _style_existing_text(label, 10, theme["muted"], True)

    title = _text_box(slide, 0.75, 1.35, 11.7, 1.55)
    _add_run(title.text_frame.paragraphs[0], deck.title, 44, theme["primary"], True)

    subtitle = _text_box(slide, 0.78, 3.12, 10.9, 0.76)
    _add_run(subtitle.text_frame.paragraphs[0], deck.subtitle, 19, theme["text"])

    footer = _text_box(slide, 0.8, 5.72, 4.8, 0.35, "Generated editable PowerPoint deck")
    _style_existing_text(footer, 10, theme["muted"])


def _add_agenda_slide(presentation, deck: DeckPlan, theme: dict) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, theme["background"])
    _add_header(slide, "Agenda", theme)

    for index, item in enumerate(deck.slides[:8], start=1):
        row = 1.55 + ((index - 1) * 0.58)
        _shape(slide, MSO_SHAPE.OVAL, 0.82, row + 0.03, 0.3, 0.3, theme["accent"])
        num = _text_box(slide, 0.82, row + 0.055, 0.3, 0.16, str(index))
        _style_existing_text(num, 7, theme["surface"], True, PP_ALIGN.CENTER)
        text = _text_box(slide, 1.28, row - 0.01, 10.7, 0.42, item.title)
        _style_existing_text(text, 18, theme["text"])


def _draw_bullet_list(slide, bullets: list[str], left, top, width, height, theme: dict, size=17) -> None:
    content = _text_box(slide, left, top, width, height)
    for index, bullet in enumerate(bullets[:5]):
        paragraph = content.text_frame.paragraphs[0] if index == 0 else content.text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(9)
        paragraph.font.size = Pt(size if len(bullet) < 115 else max(size - 3, 13))
        paragraph.font.color.rgb = theme["text"]


def _add_visual_panel(slide, theme: dict, visual_hint: str, left=8.55, top=1.55, width=3.95, height=4.88) -> None:
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, theme["primary"], theme["primary"])
    hint = _text_box(slide, left + 0.35, top + 0.55, width - 0.7, 1.95, visual_hint or "Supporting visual")
    _style_existing_text(hint, 18, theme["surface"], True, PP_ALIGN.CENTER)
    _shape(slide, MSO_SHAPE.RECTANGLE, left + 1.05, top + 2.95, width - 2.1, 0.08, theme["accent"])


def _add_bullets_layout(slide, slide_plan: SlidePlan, theme: dict) -> None:
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 1.55, 7.45, 4.88, theme["surface"], theme["surface"])
    _draw_bullet_list(slide, slide_plan.bullets, 1.05, 1.88, 6.8, 4.25, theme, 18)
    _add_visual_panel(slide, theme, slide_plan.visual_hint)


def _add_two_column_layout(slide, slide_plan: SlidePlan, theme: dict) -> None:
    left_items = slide_plan.bullets[::2] or slide_plan.bullets[:2]
    right_items = slide_plan.bullets[1::2] or slide_plan.bullets[2:]
    columns = [(0.85, "Priority"), (6.9, "Action")]
    for idx, (left, heading) in enumerate(columns):
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 1.62, 5.42, 4.76, theme["surface"], theme["surface"])
        marker = _text_box(slide, left + 0.35, 1.95, 4.7, 0.35, heading if idx == 0 else (slide_plan.visual_hint or heading))
        _style_existing_text(marker, 13, theme["primary"], True)
        _draw_bullet_list(slide, left_items if idx == 0 else right_items, left + 0.35, 2.52, 4.7, 3.35, theme, 16)


def _add_timeline_layout(slide, slide_plan: SlidePlan, theme: dict) -> None:
    bullets = slide_plan.bullets[:5]
    _shape(slide, MSO_SHAPE.RECTANGLE, 1.22, 3.3, 10.8, 0.06, theme["accent"])
    for index, bullet in enumerate(bullets):
        left = 1.0 + index * (10.5 / max(len(bullets), 1))
        _shape(slide, MSO_SHAPE.OVAL, left, 3.08, 0.48, 0.48, theme["primary"])
        num = _text_box(slide, left, 3.2, 0.48, 0.18, str(index + 1))
        _style_existing_text(num, 8, theme["surface"], True, PP_ALIGN.CENTER)
        card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left - 0.3, 3.82, 2.15, 1.45, theme["surface"], theme["surface"])
        card.shadow.inherit = False
        text = _text_box(slide, left - 0.08, 4.02, 1.7, 0.9, bullet)
        _style_existing_text(text, 11, theme["text"], False, PP_ALIGN.CENTER)


def _add_metrics_layout(slide, slide_plan: SlidePlan, theme: dict) -> None:
    bullets = slide_plan.bullets[:4]
    for index, bullet in enumerate(bullets):
        col = index % 2
        row = index // 2
        left = 0.9 + col * 6.05
        top = 1.68 + row * 2.33
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, 5.35, 1.78, theme["surface"], theme["surface"])
        number = _text_box(slide, left + 0.35, top + 0.26, 1.0, 0.45, f"0{index + 1}")
        _style_existing_text(number, 22, theme["accent"], True)
        text = _text_box(slide, left + 1.35, top + 0.28, 3.55, 0.95, bullet)
        _style_existing_text(text, 15, theme["text"], True)


def _add_quote_layout(slide, slide_plan: SlidePlan, theme: dict) -> None:
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 1.65, 11.65, 3.0, theme["surface"], theme["surface"])
    quote = slide_plan.bullets[0] if slide_plan.bullets else slide_plan.title
    quote_box = _text_box(slide, 1.28, 2.05, 10.65, 1.15, quote)
    _style_existing_text(quote_box, 25 if len(quote) < 110 else 20, theme["primary"], True, PP_ALIGN.CENTER)
    _shape(slide, MSO_SHAPE.RECTANGLE, 5.42, 3.52, 2.45, 0.08, theme["accent"])
    _draw_bullet_list(slide, slide_plan.bullets[1:] or slide_plan.bullets, 1.2, 5.0, 10.9, 1.05, theme, 14)


def _add_closing_layout(slide, slide_plan: SlidePlan, theme: dict) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.2, SLIDE_H, theme["accent"])
    message = slide_plan.bullets[0] if slide_plan.bullets else "Move forward with a clear next step."
    title = _text_box(slide, 1.0, 1.85, 11.4, 1.1, message)
    _style_existing_text(title, 28 if len(message) < 115 else 22, theme["primary"], True, PP_ALIGN.CENTER)
    _draw_bullet_list(slide, slide_plan.bullets[1:] or ["Confirm owners", "Set timeline", "Measure results"], 2.3, 3.55, 8.8, 1.8, theme, 17)


def _add_content_slide(presentation, slide_plan: SlidePlan, theme: dict, number: int, total: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, theme["background"])
    _add_header(slide, slide_plan.title, theme)

    layout = slide_plan.layout
    if layout == "two_column":
        _add_two_column_layout(slide, slide_plan, theme)
    elif layout == "timeline":
        _add_timeline_layout(slide, slide_plan, theme)
    elif layout == "metrics":
        _add_metrics_layout(slide, slide_plan, theme)
    elif layout == "quote":
        _add_quote_layout(slide, slide_plan, theme)
    elif layout == "closing":
        _add_closing_layout(slide, slide_plan, theme)
    else:
        _add_bullets_layout(slide, slide_plan, theme)

    _add_notes_and_footer(slide, slide_plan, theme, number, total)


def build_pptx(deck: DeckPlan, theme_name: str = "modern", deck_type: str = "general") -> BytesIO:
    theme = THEMES.get(theme_name.lower(), THEMES["modern"])
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    total = len(deck.slides) + 2

    _add_title_slide(presentation, deck, theme, deck_type)
    _add_agenda_slide(presentation, deck, theme)

    for index, slide_plan in enumerate(deck.slides, start=3):
        _add_content_slide(presentation, slide_plan, theme, index - 1, total - 1)

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output